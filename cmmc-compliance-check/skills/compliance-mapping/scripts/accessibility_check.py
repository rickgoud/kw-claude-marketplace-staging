#!/usr/bin/env python3
"""Real structural accessibility check for Signal D
(section-508-compliance-check, wcag-compliance-check).

Replaces the previous non-functional heuristic: content-extract only shells
out to pdftotext/pandoc, and none of pdftotext/pandoc/a plain zip-read can see
PDF tag structure, document language metadata, slide/document alt-text, or
HTML heading structure at all. This script actually executes what Signal D
has always described -- using pikepdf (PDF), python-docx (DOCX), python-pptx
(PPTX), and the stdlib html.parser (HTML), all pure-Python/C-extension
libraries, pip-installable (except html.parser, which is stdlib), no
Java/veraPDF/browser dependency required.

This is still a structural heuristic, not a full WCAG or PDF/UA conformance
test -- real conformance testing needs interactive tools (screen reader,
keyboard-navigation testing) and human review, which this cannot replace.
Say so every time this script's output is used in a report. What this DOES
check, for real, that a text-extraction-only heuristic could not:

  PDF (.pdf):
  - tagged_pdf:        PDF has MarkInfo.Marked=true AND a StructTreeRoot
                        (screen readers need both to determine reading
                        order and semantic structure)
  - language_metadata: document-level /Lang is set
  - title_metadata:    document Title metadata is set (not just the filename)
  - figures_total / figures_with_alt: walks the structure tree (only
                        possible if tagged) counting /Figure elements and
                        how many carry an /Alt entry

  DOCX (.docx):
  - language_metadata / title_metadata: core_properties.language / .title
  - figures_total / figures_with_alt: each inline image's docPr descr/title
                        attribute

  PPTX (.pptx):
  - language_metadata / title_metadata: core_properties.language / .title
  - slides_total / slides_with_title: each slide's title placeholder text
                        (screen readers and slide-navigation tools rely on
                        per-slide titles, not just the deck title)
  - figures_total / figures_with_alt: each picture shape's cNvPr descr/title
                        attribute (walked recursively into group shapes).
                        Caveat: some authoring tools (including python-pptx
                        itself, when used to insert an image programmatically)
                        auto-populate descr with the source filename rather
                        than leaving it blank -- a non-empty descr means the
                        attribute is present, not that it is meaningful
                        alt text. Same class of limitation as the PDF/DOCX
                        checks above: presence, not quality, is what's tested.

  HTML (.html/.htm), parsed with the stdlib html.parser -- no browser or
  rendering engine, so this cannot see content injected by JavaScript:
  - language_metadata: <html lang="..."> (or xml:lang) attribute is set
  - title_metadata:    <title> element has non-empty text
  - heading_hierarchy_ok: heading levels (h1-h6) in document order never
                        skip a level going deeper (e.g. h1 -> h3 with no h2)
  - has_h1:            at least one <h1> element is present
  - figures_total / figures_with_alt: each <img> counted, "has alt" means
                        the alt attribute is present at all (including
                        alt="" for intentionally-decorative images, which is
                        valid under WCAG); a missing alt attribute entirely
                        is the failure this flags

Only .pdf, .docx, .pptx, .html, and .htm are supported today -- other formats
are reported as skipped, not silently ignored, so a report can say plainly
what wasn't checked.

Usage:
    python3 accessibility_check.py <path-to-file>
Prints a JSON result to stdout. Install deps once per environment:
    pip install pikepdf python-docx python-pptx --break-system-packages
"""
import sys
import json
import os


def check_pdf(path):
    import pikepdf

    result = {"file": path, "type": "pdf", "issues": [], "checks": {}}
    try:
        pdf = pikepdf.open(path)
    except Exception as e:
        result["error"] = "could not open PDF: %s" % e
        return result

    root = pdf.Root

    # -- Tagged?
    marked = False
    try:
        marked = bool(root.MarkInfo.Marked)
    except Exception:
        marked = False
    has_struct_tree = "/StructTreeRoot" in root
    tagged = marked and has_struct_tree
    result["checks"]["tagged_pdf"] = tagged
    if not tagged:
        result["issues"].append(
            "Not a tagged PDF (no MarkInfo.Marked + StructTreeRoot) -- screen "
            "readers cannot reliably determine reading order or semantic structure."
        )

    # -- Language
    lang = None
    try:
        if "/Lang" in root:
            lang = str(root.Lang)
    except Exception:
        lang = None
    result["checks"]["language_metadata"] = bool(lang)
    result["checks"]["language_value"] = lang
    if not lang:
        result["issues"].append(
            "No document-level /Lang entry -- assistive technology cannot "
            "determine the document's language."
        )

    # -- Title
    title = None
    try:
        if pdf.docinfo is not None and "/Title" in pdf.docinfo:
            title = str(pdf.docinfo["/Title"])
    except Exception:
        title = None
    result["checks"]["title_metadata"] = bool(title)
    if not title:
        result["issues"].append(
            "No document Title metadata set -- screen readers announce the "
            "filename instead of a meaningful title."
        )

    # -- Alt text on figures (only walkable if tagged)
    figures_total = 0
    figures_with_alt = 0
    if has_struct_tree:
        try:
            import pikepdf as pk

            def walk(node, depth=0):
                nonlocal figures_total, figures_with_alt
                if depth > 40 or node is None:
                    return
                try:
                    s_type = node.get("/S", None) if hasattr(node, "get") else None
                except Exception:
                    s_type = None
                if s_type is not None and str(s_type) == "/Figure":
                    figures_total += 1
                    try:
                        if "/Alt" in node:
                            figures_with_alt += 1
                    except Exception:
                        pass
                try:
                    kids = node.get("/K", None) if hasattr(node, "get") else None
                except Exception:
                    kids = None
                if kids is None:
                    return
                if isinstance(kids, pk.Array):
                    for k in kids:
                        walk(k, depth + 1)
                else:
                    walk(kids, depth + 1)

            struct_root = root.StructTreeRoot
            k = struct_root.get("/K", None) if "/K" in struct_root else None
            if isinstance(k, pk.Array):
                for item in k:
                    walk(item)
            elif k is not None:
                walk(k)
        except Exception as e:
            result["issues"].append(
                "Could not fully walk structure tree for alt-text check: %s" % e
            )

    result["checks"]["figures_total"] = figures_total
    result["checks"]["figures_with_alt"] = figures_with_alt
    if figures_total > 0 and figures_with_alt < figures_total:
        result["issues"].append(
            "%d of %d tagged figures have no /Alt text."
            % (figures_total - figures_with_alt, figures_total)
        )

    pdf.close()
    return result


def check_docx(path):
    import docx
    from docx.oxml.ns import qn

    result = {"file": path, "type": "docx", "issues": [], "checks": {}}
    try:
        d = docx.Document(path)
    except Exception as e:
        result["error"] = "could not open DOCX: %s" % e
        return result

    title = d.core_properties.title
    result["checks"]["title_metadata"] = bool(title)
    if not title:
        result["issues"].append("No document Title set in core properties.")

    lang = d.core_properties.language
    result["checks"]["language_metadata"] = bool(lang)
    if not lang:
        result["issues"].append("No document language set in core properties.")

    figures_total = 0
    figures_with_alt = 0
    try:
        for shape in d.inline_shapes:
            figures_total += 1
            doc_pr = shape._inline.find(qn("wp:docPr"))
            descr = doc_pr.get("descr") if doc_pr is not None else None
            title_attr = doc_pr.get("title") if doc_pr is not None else None
            if descr or title_attr:
                figures_with_alt += 1
    except Exception as e:
        result["issues"].append("Could not check inline image alt text: %s" % e)

    result["checks"]["figures_total"] = figures_total
    result["checks"]["figures_with_alt"] = figures_with_alt
    if figures_total > 0 and figures_with_alt < figures_total:
        result["issues"].append(
            "%d of %d images have no alt text/description set."
            % (figures_total - figures_with_alt, figures_total)
        )

    return result


def check_pptx(path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    result = {"file": path, "type": "pptx", "issues": [], "checks": {}}
    try:
        prs = Presentation(path)
    except Exception as e:
        result["error"] = "could not open PPTX: %s" % e
        return result

    title = prs.core_properties.title
    result["checks"]["title_metadata"] = bool(title)
    if not title:
        result["issues"].append("No document Title set in core properties.")

    lang = prs.core_properties.language
    result["checks"]["language_metadata"] = bool(lang)
    if not lang:
        result["issues"].append("No document language set in core properties.")

    def iter_shapes(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub in iter_shapes(shape.shapes):
                        yield sub
                except Exception:
                    pass

    slides_total = 0
    slides_with_title = 0
    figures_total = 0
    figures_with_alt = 0

    try:
        for slide in prs.slides:
            slides_total += 1
            try:
                t = slide.shapes.title
                if t is not None and t.has_text_frame and t.text_frame.text.strip():
                    slides_with_title += 1
            except Exception:
                pass

            try:
                for shape in iter_shapes(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        figures_total += 1
                        descr = None
                        title_attr = None
                        try:
                            cnv_pr = shape._element.find(".//" + qn("p:cNvPr"))
                            if cnv_pr is not None:
                                descr = cnv_pr.get("descr")
                                title_attr = cnv_pr.get("title")
                        except Exception:
                            pass
                        if descr or title_attr:
                            figures_with_alt += 1
            except Exception as e:
                result["issues"].append(
                    "Could not check slide shapes for alt text: %s" % e
                )
    except Exception as e:
        result["issues"].append("Could not walk slides: %s" % e)

    result["checks"]["slides_total"] = slides_total
    result["checks"]["slides_with_title"] = slides_with_title
    if slides_total > 0 and slides_with_title < slides_total:
        result["issues"].append(
            "%d of %d slides have no title placeholder text -- screen readers "
            "and slide-navigation tools rely on per-slide titles."
            % (slides_total - slides_with_title, slides_total)
        )

    result["checks"]["figures_total"] = figures_total
    result["checks"]["figures_with_alt"] = figures_with_alt
    if figures_total > 0 and figures_with_alt < figures_total:
        result["issues"].append(
            "%d of %d images have no alt text/description set."
            % (figures_total - figures_with_alt, figures_total)
        )

    return result


class _A11yHTMLParser:
    """Minimal, dependency-free HTML structure walker built on the stdlib
    html.parser. No rendering engine -- content injected by JavaScript after
    load is invisible to this, same limitation as every other check here
    that reads static file bytes rather than executing the document."""

    def __init__(self):
        from html.parser import HTMLParser

        outer = self

        class _Parser(HTMLParser):
            def __init__(self):
                super().__init__(convert_charrefs=True)
                self.lang = None
                self.title = ""
                self._in_title = False
                self.headings = []  # list of ints, document order
                self.images = []  # list of bool, has_alt

            def handle_starttag(self, tag, attrs):
                tag = tag.lower()
                attrs_d = {k.lower(): v for k, v in attrs}
                if tag == "html":
                    self.lang = attrs_d.get("lang") or attrs_d.get("xml:lang")
                elif tag == "title":
                    self._in_title = True
                elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                    self.headings.append(int(tag[1]))
                elif tag == "img":
                    self.images.append("alt" in attrs_d)

            def handle_startendtag(self, tag, attrs):
                # self-closing tags, e.g. <img ... />
                self.handle_starttag(tag, attrs)

            def handle_data(self, data):
                if self._in_title:
                    self.title += data

            def handle_endtag(self, tag):
                if tag.lower() == "title":
                    self._in_title = False

        self._parser = _Parser()

    def feed(self, text):
        self._parser.feed(text)

    @property
    def lang(self):
        return self._parser.lang

    @property
    def title(self):
        return self._parser.title

    @property
    def headings(self):
        return self._parser.headings

    @property
    def images(self):
        return self._parser.images


def check_html(path):
    result = {"file": path, "type": "html", "issues": [], "checks": {}}
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
    except Exception as e:
        result["error"] = "could not open HTML: %s" % e
        return result

    parser = _A11yHTMLParser()
    try:
        parser.feed(text)
    except Exception as e:
        result["issues"].append("Could not fully parse HTML: %s" % e)

    lang = parser.lang
    result["checks"]["language_metadata"] = bool(lang)
    result["checks"]["language_value"] = lang
    if not lang:
        result["issues"].append(
            "No <html lang=\"...\"> attribute -- assistive technology cannot "
            "determine the page's language."
        )

    title = (parser.title or "").strip()
    result["checks"]["title_metadata"] = bool(title)
    if not title:
        result["issues"].append(
            "No non-empty <title> element -- screen readers announce the URL "
            "instead of a meaningful title."
        )

    headings = parser.headings
    has_h1 = 1 in headings
    result["checks"]["has_h1"] = has_h1
    if headings and not has_h1:
        result["issues"].append("No <h1> element found on the page.")

    skips = []
    prev_level = 0
    for lvl in headings:
        if prev_level and lvl > prev_level + 1:
            skips.append((prev_level, lvl))
        prev_level = lvl
    result["checks"]["heading_hierarchy_ok"] = len(skips) == 0
    if skips:
        result["issues"].append(
            "Heading level(s) skipped: "
            + ", ".join("h%d -> h%d" % (a, b) for a, b in skips)
            + " -- screen reader users navigating by heading level will miss "
            "the skipped level(s)."
        )

    images = parser.images
    figures_total = len(images)
    figures_with_alt = sum(1 for has_alt in images if has_alt)
    result["checks"]["figures_total"] = figures_total
    result["checks"]["figures_with_alt"] = figures_with_alt
    if figures_total > 0 and figures_with_alt < figures_total:
        result["issues"].append(
            "%d of %d <img> elements have no alt attribute at all."
            % (figures_total - figures_with_alt, figures_total)
        )

    return result


def check_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return check_pdf(path)
    elif ext == ".docx":
        return check_docx(path)
    elif ext == ".pptx":
        return check_pptx(path)
    elif ext in (".html", ".htm"):
        return check_html(path)
    else:
        return {
            "file": path,
            "type": ext.lstrip("."),
            "checks": {},
            "issues": [],
            "skipped": (
                "unsupported file type for structural accessibility check "
                "(only .pdf, .docx, .pptx, .html, and .htm are checked today)"
            ),
        }


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("usage: accessibility_check.py <path-to-file>", file=sys.stderr)
        sys.exit(1)
    print(json.dumps(check_file(sys.argv[1]), indent=2))
